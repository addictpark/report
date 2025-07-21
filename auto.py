import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# 파일명 및 시트/테이블명 매핑
excel_file = 'autotest.xlsx'
ppt_file = '자동화.pptx'
output_file = 'output_report.pptx'

sheet_to_table = {
    # 이용 인원
    '서비스 이용 인원': 'service1',
    '상담유형별 이용 인원': 'type_people',
    '성별 이용 인원': 'sex_people',
    '연령별 이용 인원': 'age_people',
    # 이용 횟수
    '서비스 이용 횟수': 'service2',
    '상담유형별 이용 횟수': 'type_case',
    '성별 이용 횟수': 'sex_case',
    '연령별 이용 횟수': 'age_case',
}

# 데이터프레임 읽기
df_dict = {sheet: pd.read_excel(excel_file, sheet_name=sheet) for sheet in sheet_to_table.keys()}

# 폰트/크기
font_name = '나눔스퀘어_ac'
font_size = Pt(10)
MIN_TABLE_ROWS = 2

def format_value(val):
    if pd.isnull(val):
        return '-'
    if isinstance(val, float) and val.is_integer():
        return str(int(val))
    return str(val)

# ---- 이용 인원 표: 누계, 실계 ----
def fill_table_people(table, df):
    min_row = MIN_TABLE_ROWS
    n_table_data = len(table.rows) - min_row - 2
    n_excel_data = len(df) - 2
    n_cols = len(table.columns)

    for i in range(n_table_data):
        if i < n_excel_data:
            row = df.iloc[i]
            for j, value in enumerate(row):
                cell = table.cell(i + min_row, j)
                cell.text = format_value(value)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = font_size
        else:
            for j in range(n_cols):
                cell = table.cell(i + min_row, j)
                cell.text = '-'
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = font_size

    # 누계, 실계
    for j, value in enumerate(df.iloc[-2]):
        cell = table.cell(len(table.rows) - 2, j)
        cell.text = format_value(value)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size

    for j, value in enumerate(df.iloc[-1]):
        cell = table.cell(len(table.rows) - 1, j)
        cell.text = format_value(value)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size

# ---- 이용 횟수 표: 누계 ----
def fill_table_count(table, df):
    min_row = MIN_TABLE_ROWS
    n_table_data = len(table.rows) - min_row - 1
    n_excel_data = len(df) - 1
    n_cols = len(table.columns)

    for i in range(n_table_data):
        if i < n_excel_data:
            row = df.iloc[i]
            for j, value in enumerate(row):
                cell = table.cell(i + min_row, j)
                cell.text = format_value(value)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = font_size
        else:
            for j in range(n_cols):
                cell = table.cell(i + min_row, j)
                cell.text = '-'
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = font_size

    # 누계
    for j, value in enumerate(df.iloc[-1]):
        cell = table.cell(len(table.rows) - 1, j)
        cell.text = format_value(value)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size

# ---- PPT에 데이터 입력 ----
prs = Presentation(ppt_file)
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_table:
            for sheet, tbl_name in sheet_to_table.items():
                if shape.name == tbl_name:
                    if '인원' in sheet:
                        fill_table_people(shape.table, df_dict[sheet])
                    elif '횟수' in sheet:
                        fill_table_count(shape.table, df_dict[sheet])

prs.save(output_file)