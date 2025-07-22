import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import io

# Streamlit UI
st.title("📊 통합형 PPT 자동 보고서 생성기")
st.markdown("PPT 템플릿과 엑셀을 업로드하고 템플릿 유형을 선택하면 자동으로 보고서를 생성합니다.")

# 템플릿 선택
template_option = st.selectbox(
    "🗂️ 사용할 템플릿 양식을 선택하세요:",
    ["세아 버전", "한국은행 버전"]
)

uploaded_ppt = st.file_uploader("📁 PPT 템플릿 파일(PPTX)을 업로드하세요", type=["pptx"])
uploaded_excel = st.file_uploader("📁 엑셀 파일(XLSX)을 업로드하세요", type=["xlsx"])

# 시트와 도형 이름 매핑
if template_option == "세아 버전":
    table_map = {
        '서비스 이용 인원': 'service1',
        '상담유형별 이용 인원': 'type_people',
        '성별 이용 인원': 'sex_people',
        '연령별 이용 인원': 'age_people',
        '서비스 이용 횟수': 'service2',
        '상담유형별 이용 횟수': 'type_case',
        '성별 이용 횟수': 'sex_case',
        '연령별 이용 횟수': 'age_case',
        '심리진단 이용 횟수': 'diag_case',
    }
elif template_option == "한국은행 버전":
    table_map = {
        '서비스 이용 인원': 'service1',
        '상담방법별 이용 인원': 'type_people',
        '상담유형별 이용 인원': 'subject_people',
        '성별 이용 인원': 'sex_people',
        '연령별 이용 인원': 'age_people',
        '소속별 이용 인원': 'group_people',
        '직급별 이용 인원': 'class_people',
        '서비스 이용 횟수': 'service2',
        '상담방법별 이용 횟수': 'type_case',
        '상담유형별 이용 횟수': 'subject_case',
        '성별 이용 횟수': 'sex_case',
        '연령별 이용 횟수': 'age_case',
        '소속별 이용 횟수': 'group_case',
        '직급별 이용 횟수': 'class_case',
        '심리진단 이용 횟수': 'diag_case',
    }

font_name = '나눔스퀘어_ac'
font_size = Pt(10)
MIN_TABLE_ROWS = 2

def format_value(val):
    if pd.isnull(val):
        return '-'
    try:
        if float(val) == 0:
            return '-'
    except:
        pass
    if isinstance(val, float):
        return str(int(val)) if val.is_integer() else str(val)
    return str(val)

def fill_table_people(table, df):
    min_row = MIN_TABLE_ROWS
    n_data = len(table.rows) - min_row - 2
    for i in range(n_data):
        row = df.iloc[i] if i < len(df) - 2 else ['-'] * len(table.columns)
        for j, val in enumerate(row):
            set_cell(table.cell(i + min_row, j), format_value(val))
    for j, val in enumerate(df.iloc[-2]):
        set_cell(table.cell(len(table.rows) - 2, j), format_value(val))
    for j, val in enumerate(df.iloc[-1]):
        set_cell(table.cell(len(table.rows) - 1, j), format_value(val))

def fill_table_count(table, df):
    min_row = MIN_TABLE_ROWS
    n_data = len(table.rows) - min_row - 1
    for i in range(n_data):
        row = df.iloc[i] if i < len(df) - 1 else ['-'] * len(table.columns)
        for j, val in enumerate(row):
            set_cell(table.cell(i + min_row, j), format_value(val))
    for j, val in enumerate(df.iloc[-1]):
        set_cell(table.cell(len(table.rows) - 1, j), format_value(val))

def set_cell(cell, text):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = font_size

if uploaded_excel and uploaded_ppt:
    df_dict = {sheet: pd.read_excel(uploaded_excel, sheet_name=sheet) for sheet in table_map.keys()}
    prs = Presentation(uploaded_ppt)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for sheet, shape_name in table_map.items():
                    if shape.name == shape_name:
                        df = df_dict[sheet]
                        if '인원' in sheet:
                            fill_table_people(shape.table, df)
                        elif '횟수' in sheet:
                            fill_table_count(shape.table, df)

    output = io.BytesIO()
    prs.save(output)
    st.success("✅ 보고서 생성이 완료되었습니다!")

    st.download_button(
        label="📥 자동화 보고서 다운로드",
        data=output.getvalue(),
        file_name="output_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
