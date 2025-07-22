import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import io

# 설정
font_name = '나눔스퀘어_ac'
font_size = Pt(10)
MIN_TABLE_ROWS = 2

# 파일 업로드
st.title("📊 PPT 자동화 보고서 생성기")
st.markdown("업로드한 엑셀 및 PPT 템플릿을 기반으로 자동화된 보고서를 생성합니다.")

uploaded_excel = st.file_uploader("📁 엑셀 파일을 업로드하세요", type=["xlsx"])
uploaded_ppt = st.file_uploader("📁 PPT 템플릿 파일을 업로드하세요", type=["pptx"])

# 시트-도형명 매핑
sheet_to_table = {
    '서비스 이용 인원': 'service1',
    '상담유형별 이용 인원': 'type_people',
    '성별 이용 인원': 'sex_people',
    '연령별 이용 인원': 'age_people',
    '서비스 이용 횟수': 'service2',
    '상담유형별 이용 횟수': 'type_case',
    '성별 이용 횟수': 'sex_case',
    '연령별 이용 횟수': 'age_case',
    '심리진단 이용 횟수': 'diag_case',
}

def format_value(val):
    if pd.isnull(val):
        return '-'
    try:
        if float(val) == 0:
            return '-'
    except:
        pass
    if isinstance(val, float):
        return str(int(val)) if val.is_integer() else str(val)
    return str(val)

def fill_table_people(table, df):
    min_row = MIN_TABLE_ROWS
    n_table_data = len(table.rows) - min_row - 2
    n_excel_data = len(df) - 2
    n_cols = len(table.columns)

    for i in range(n_table_data):
        row = df.iloc[i] if i < n_excel_data else ['-'] * n_cols
        for j, value in enumerate(row):
            cell = table.cell(i + min_row, j)
            cell.text = format_value(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = font_size

    # 누계 / 실계
    for offset in [-2, -1]:
        for j, value in enumerate(df.iloc[offset]):
            cell = table.cell(len(table.rows) + offset, j)
            cell.text = format_value(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = font_size

def fill_table_count(table, df):
    min_row = MIN_TABLE_ROWS
    n_table_data = len(table.rows) - min_row - 1
    n_excel_data = len(df) - 1
    n_cols = len(table.columns)

    for i in range(n_table_data):
        row = df.iloc[i] if i < n_excel_data else ['-'] * n_cols
        for j, value in enumerate(row):
            cell = table.cell(i + min_row, j)
            cell.text = format_value(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = font_size

    # 누계
    for j, value in enumerate(df.iloc[-1]):
        cell = table.cell(len(table.rows) - 1, j)
        cell.text = format_value(value)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size

if uploaded_excel and uploaded_ppt:
    df_dict = {
        sheet: pd.read_excel(uploaded_excel, sheet_name=sheet)
        for sheet in sheet_to_table.keys()
    }

    prs = Presentation(uploaded_ppt)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for sheet, tbl_name in sheet_to_table.items():
                    if shape.name == tbl_name:
                        if '인원' in sheet:
                            fill_table_people(shape.table, df_dict[sheet])
                        elif '횟수' in sheet:
                            fill_table_count(shape.table, df_dict[sheet])

    output = io.BytesIO()
    prs.save(output)
    st.success("✅ 보고서 생성 완료!")

    st.download_button(
        label="📥 자동화 보고서 다운로드",
        data=output.getvalue(),
        file_name="output_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
